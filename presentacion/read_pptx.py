import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

prs = Presentation(r'presentacion\Sistema UniTime.pptx')
print(f'Total slides: {len(prs.slides)}')
print()

for i, slide in enumerate(prs.slides):
    print(f'===== SLIDE {i+1} =====')
    layout = slide.slide_layout.name if slide.slide_layout else 'N/A'
    print(f'Layout: {layout}')
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    print(f'  {text}')
        if shape.has_table:
            table = shape.table
            nrows = len(table.rows)
            ncols = len(table.columns)
            print(f'  [TABLE {nrows}x{ncols}]')
            for row in table.rows:
                cells = []
                for cell in row.cells:
                    cells.append(cell.text.strip())
                line = ' | '.join(cells)
                print(f'    | {line} |')
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            print(f'  [IMAGE]')
    print()
